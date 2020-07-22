import base64
import csv
import datetime
import json

from docx import Document
from docx.shared import Inches
from PIL import ImageColor

from main import app
from core import action_hub, get_output_file_name, get_temp_file_name, get_sdk_for_schedule, get_sdk_all_access, send_email
from api_types import ActionDefinition, ActionList, ActionFormField, ActionRequest, ActionForm

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION






slug = 'look_to_powerpoint'

icon_data_uri = 'data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAgAAAAIACAYAAAD0eNT6AAAACXBIWXMAAA7CAAAOwgEVKEqAAAAKT2lDQ1BQaG90b3Nob3AgSUNDIHByb2ZpbGUAAHjanVNnVFPpFj333vRCS4iAlEtvUhUIIFJCi4AUkSYqIQkQSoghodkVUcERRUUEG8igiAOOjoCMFVEsDIoK2AfkIaKOg6OIisr74Xuja9a89+bN/rXXPues852zzwfACAyWSDNRNYAMqUIeEeCDx8TG4eQuQIEKJHAAEAizZCFz/SMBAPh+PDwrIsAHvgABeNMLCADATZvAMByH/w/qQplcAYCEAcB0kThLCIAUAEB6jkKmAEBGAYCdmCZTAKAEAGDLY2LjAFAtAGAnf+bTAICd+Jl7AQBblCEVAaCRACATZYhEAGg7AKzPVopFAFgwABRmS8Q5ANgtADBJV2ZIALC3AMDOEAuyAAgMADBRiIUpAAR7AGDIIyN4AISZABRG8lc88SuuEOcqAAB4mbI8uSQ5RYFbCC1xB1dXLh4ozkkXKxQ2YQJhmkAuwnmZGTKBNA/g88wAAKCRFRHgg/P9eM4Ors7ONo62Dl8t6r8G/yJiYuP+5c+rcEAAAOF0ftH+LC+zGoA7BoBt/qIl7gRoXgugdfeLZrIPQLUAoOnaV/Nw+H48PEWhkLnZ2eXk5NhKxEJbYcpXff5nwl/AV/1s+X48/Pf14L7iJIEyXYFHBPjgwsz0TKUcz5IJhGLc5o9H/LcL//wd0yLESWK5WCoU41EScY5EmozzMqUiiUKSKcUl0v9k4t8s+wM+3zUAsGo+AXuRLahdYwP2SycQWHTA4vcAAPK7b8HUKAgDgGiD4c93/+8//UegJQCAZkmScQAAXkQkLlTKsz/HCAAARKCBKrBBG/TBGCzABhzBBdzBC/xgNoRCJMTCQhBCCmSAHHJgKayCQiiGzbAdKmAv1EAdNMBRaIaTcA4uwlW4Dj1wD/phCJ7BKLyBCQRByAgTYSHaiAFiilgjjggXmYX4IcFIBBKLJCDJiBRRIkuRNUgxUopUIFVIHfI9cgI5h1xGupE7yAAygvyGvEcxlIGyUT3UDLVDuag3GoRGogvQZHQxmo8WoJvQcrQaPYw2oefQq2gP2o8+Q8cwwOgYBzPEbDAuxsNCsTgsCZNjy7EirAyrxhqwVqwDu4n1Y8+xdwQSgUXACTYEd0IgYR5BSFhMWE7YSKggHCQ0EdoJNwkDhFHCJyKTqEu0JroR+cQYYjIxh1hILCPWEo8TLxB7iEPENyQSiUMyJ7mQAkmxpFTSEtJG0m5SI+ksqZs0SBojk8naZGuyBzmULCAryIXkneTD5DPkG+Qh8lsKnWJAcaT4U+IoUspqShnlEOU05QZlmDJBVaOaUt2ooVQRNY9aQq2htlKvUYeoEzR1mjnNgxZJS6WtopXTGmgXaPdpr+h0uhHdlR5Ol9BX0svpR+iX6AP0dwwNhhWDx4hnKBmbGAcYZxl3GK+YTKYZ04sZx1QwNzHrmOeZD5lvVVgqtip8FZHKCpVKlSaVGyovVKmqpqreqgtV81XLVI+pXlN9rkZVM1PjqQnUlqtVqp1Q61MbU2epO6iHqmeob1Q/pH5Z/YkGWcNMw09DpFGgsV/jvMYgC2MZs3gsIWsNq4Z1gTXEJrHN2Xx2KruY/R27iz2qqaE5QzNKM1ezUvOUZj8H45hx+Jx0TgnnKKeX836K3hTvKeIpG6Y0TLkxZVxrqpaXllirSKtRq0frvTau7aedpr1Fu1n7gQ5Bx0onXCdHZ4/OBZ3nU9lT3acKpxZNPTr1ri6qa6UbobtEd79up+6Ynr5egJ5Mb6feeb3n+hx9L/1U/W36p/VHDFgGswwkBtsMzhg8xTVxbzwdL8fb8VFDXcNAQ6VhlWGX4YSRudE8o9VGjUYPjGnGXOMk423GbcajJgYmISZLTepN7ppSTbmmKaY7TDtMx83MzaLN1pk1mz0x1zLnm+eb15vft2BaeFostqi2uGVJsuRaplnutrxuhVo5WaVYVVpds0atna0l1rutu6cRp7lOk06rntZnw7Dxtsm2qbcZsOXYBtuutm22fWFnYhdnt8Wuw+6TvZN9un2N/T0HDYfZDqsdWh1+c7RyFDpWOt6azpzuP33F9JbpL2dYzxDP2DPjthPLKcRpnVOb00dnF2e5c4PziIuJS4LLLpc+Lpsbxt3IveRKdPVxXeF60vWdm7Obwu2o26/uNu5p7ofcn8w0nymeWTNz0MPIQ+BR5dE/C5+VMGvfrH5PQ0+BZ7XnIy9jL5FXrdewt6V3qvdh7xc+9j5yn+M+4zw33jLeWV/MN8C3yLfLT8Nvnl+F30N/I/9k/3r/0QCngCUBZwOJgUGBWwL7+Hp8Ib+OPzrbZfay2e1BjKC5QRVBj4KtguXBrSFoyOyQrSH355jOkc5pDoVQfujW0Adh5mGLw34MJ4WHhVeGP45wiFga0TGXNXfR3ENz30T6RJZE3ptnMU85ry1KNSo+qi5qPNo3ujS6P8YuZlnM1VidWElsSxw5LiquNm5svt/87fOH4p3iC+N7F5gvyF1weaHOwvSFpxapLhIsOpZATIhOOJTwQRAqqBaMJfITdyWOCnnCHcJnIi/RNtGI2ENcKh5O8kgqTXqS7JG8NXkkxTOlLOW5hCepkLxMDUzdmzqeFpp2IG0yPTq9MYOSkZBxQqohTZO2Z+pn5mZ2y6xlhbL+xW6Lty8elQfJa7OQrAVZLQq2QqboVFoo1yoHsmdlV2a/zYnKOZarnivN7cyzytuQN5zvn//tEsIS4ZK2pYZLVy0dWOa9rGo5sjxxedsK4xUFK4ZWBqw8uIq2Km3VT6vtV5eufr0mek1rgV7ByoLBtQFr6wtVCuWFfevc1+1dT1gvWd+1YfqGnRs+FYmKrhTbF5cVf9go3HjlG4dvyr+Z3JS0qavEuWTPZtJm6ebeLZ5bDpaql+aXDm4N2dq0Dd9WtO319kXbL5fNKNu7g7ZDuaO/PLi8ZafJzs07P1SkVPRU+lQ27tLdtWHX+G7R7ht7vPY07NXbW7z3/T7JvttVAVVN1WbVZftJ+7P3P66Jqun4lvttXa1ObXHtxwPSA/0HIw6217nU1R3SPVRSj9Yr60cOxx++/p3vdy0NNg1VjZzG4iNwRHnk6fcJ3/ceDTradox7rOEH0x92HWcdL2pCmvKaRptTmvtbYlu6T8w+0dbq3nr8R9sfD5w0PFl5SvNUyWna6YLTk2fyz4ydlZ19fi753GDborZ752PO32oPb++6EHTh0kX/i+c7vDvOXPK4dPKy2+UTV7hXmq86X23qdOo8/pPTT8e7nLuarrlca7nuer21e2b36RueN87d9L158Rb/1tWeOT3dvfN6b/fF9/XfFt1+cif9zsu72Xcn7q28T7xf9EDtQdlD3YfVP1v+3Njv3H9qwHeg89HcR/cGhYPP/pH1jw9DBY+Zj8uGDYbrnjg+OTniP3L96fynQ89kzyaeF/6i/suuFxYvfvjV69fO0ZjRoZfyl5O/bXyl/erA6xmv28bCxh6+yXgzMV70VvvtwXfcdx3vo98PT+R8IH8o/2j5sfVT0Kf7kxmTk/8EA5jz/GMzLdsAAAAgY0hSTQAAeiUAAICDAAD5/wAAgOkAAHUwAADqYAAAOpgAABdvkl/FRgAAK11JREFUeNrs3XmUXFdh5/HfvW+rpau6q7urqquqN1leICwmYMAsCTsEAsOaEMjCJDBkyDaTwwmTgZxJMnNmJieZCZkTkgmZk4RJIMNMIJvDIRBMDMEYO8E2Jt6NbCNZq9Vqtbqra3nv3vnjlWUZZMDQkrq6v59jobYsJJ37unW/77737jM3v+QSLwAAsKtYhgAAAAIAAAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAAgAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAEAAAAIAAAAAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAACAAAAEAAAAIAAAAAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAAAIAAAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAAgAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAEAAAAIAAAAAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAACgCEAAIAAAAAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAACAAAAAgAAAAAAEAAAAIAAAAAABAAAAxlbIEAB47Hz+j/dSlsm7VD5NZYJQtlBkeAACAMDOmO/daLJ38i6TvJfkZaJY4dS0wpmG4npL6dpJbX7lVinLJGMYN4AAADA+k73PJ3vv5NNUkpGNIpkgVDgzo6g+p6SzrLjRVmFpr6KZuTwAGi2d/Ie/0/73vkduOJAJ+OsFIAAAbNvJ/msnfRsXFJTKsuUJJXMLStpLSpYuVtJeUtyaV1CaUDBRlS2Wv/6XUyZ5J4mzf4AAALB9JnhJMiZfnTdWJk7yJfzJaUWzcyosXpRP+MuXKmm0ZIpl2SiWiZNv7bfLMsYcIAAAnO/J3nsnufyb914mDGWCUCZOFMSJwqkZRXPziustxe1FFdqLStqLiluL+Rm95cEggAAAsL3n+yyTXCafZfJZKhNGCgol2YmSgvKEgmotn+Q7y4rm5pXMzSuabSlqtLhTHwABAIzFZO+clKXyWTq6MU8KqzUF1al8CX9qRnGzo2QhX8KP6nOKZpsKKlOy3+LyPQACAMAFne19flbvUinN5F2msDqlqL2oaHZOcbOtcLqhwvyy4vaiopmmwqkZBaWyZAPGDwABAGz3if70hD+6A1/OyRbL+Zl8Z0lRs6NCZ0lRvaWo2VE0XVcwMSmbFBg/AAQAsJ0n+HxyH30zkrGBZK1MEMgmJcWtecXtJRWXL1E8uikvnKwpqE7LFooyYcRYAiAAgO070Y/uwD/jkTsbxfmNecWiglJF4UxDhcW9ittLKixfoqS1oHBqRiYuyCZcqwdAAADb1+gRO7n8+nx+Wm9k41hBZTK/Hj96rj5pLSrpLCqeW1DSWVY4XZcJuE4PgAAAtvnZvcvvwHdOPsvvvjdhJBNFCifrCqdmFbcW8jvwO4uKG21FjY6imYbCqRnGDwABAGz/yX60dD+a7I0kE8cKJiYVjLbFjZptFRYvVrJ4kaLarIJqLT/jn6gyfgAIAGC7T/SjD0YvtcvvwDdJQVFlSuHktOJGS/H8spLWkpLlixXPNBWUq7LFkmyxxBgCIACAbTvJey+vh+7Al4y1+Stqg0DGBoqm64oa7Xyy7ywp6SznL7xpdmQLJZkolgn5EgFAAADbdrL3zuXX67NMkn94//sokY0TBZXJ/BW2ncX8ZTftRYWzLcXNtoJyhTEEQAAAYzHZZ6m8y+TTVCYM89fTliYUlKsKq1OK5+aVLO7NJ/u5eYUzDYWTtbO+whYACABgu833LpOyfKLP78I3CqqTiuY6+ZvtZhqKGx0lC3tO340f1eqy5QmZgE9vACAAMCZn96MJP8vkvVNUnVZYbypuLSputPKNc+otFVqL+Vl9bZZtcQGAAMA2n+El//DHp+/AN0bGSLZQUjy3kL/CdvEixc0FJfPLCidrCmt1BeUJmShmGAGAAMB2PIs/fTafz/Cj/2BkrJWJYtmkIFssKZqdy7fFbS3m2+K2lxRUa7KFomyxmO+ZDwAgALANJ/vTm+dk8t7LmHxbXBNGCopFBZVa/irbeuuM99UvKunsUdRo5Y/mAQAIAGxjzuUvvBldo5c1+eN2UaxgYlrh5IyiZltxvaV4rpMv57cWFTfbiqbrjB8AEAAYh7N7n2X5mX2WynufL80XygqrVUUzTUXNTv6YXWdR0UxL0WxT0UxDQXWK8QMAAgBjMd87J7lMck4uTWULBYXTswqrtdELbtr5jXmdPYpm6gqnpvM98ksTDB4AEADYzmf0X/uxl5eclwkCxbON02+xi+cvUtJaUGFhj8LZlsLyhGyhKBPzvnoAIACwbSd67/3pF9zI59viytr8en0QKKzNjt5Pv6S4taikvaio2VJc78iWyvke+NycBwAEALbxZD9atn/oexNGsoWSbJLIJiXZUlnF5UsUd/YoWdijpL10+jo9r7AFAAIAYzLZ+yzL98HPsvxRu4l873tbriicmlYyt6Bk6WIlrUXF7SVFs00FxRLL9wAAAmBsGCM/HMr1NyVjFE5O58v3s02F0/XREn7+Zrtodk5RvSmbFBk3AAABMNYn/YOB4s6Spp77UkWNtpLOksLarKJaXeFkjbN6AAABsBO5XlcTT7pCrbf/O7bFBQBsKW713vZHKGDyBwAQALtvGcAxBgAAAgAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAABAAAACAAAAEAAAAAAAgAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAEAAAAODbEzIEAM4nn6WSc5IxDMa4M1bGWo4lAQAA32S+kBTVW7LlipSmDMg4h5zL5Hpdpasr8oOBTBwzKAQAAJxtxvAySUHtn3y3ypc/Q/KeMRnnw5kOlZ54UJv77tT6jdfq5Oc/Jd/flAKmFQIAAL5uCcAonKkrmq4zFjtA3Oyo9LjLNf2y1+nk5z6pwx94r3r33yMTRgzOGOAmQADn76wxS+U2uwzETuu6INTU816hpff8lgrLl8r3ewwKAQAAj5gquGFsByte/AS13/Ee2dKElGUMCAEAANgtKt99paZf9jo5VgEIAADA7mHCWNUrX6igOpk/8gkCAACwO8TNecXNeXke9SQAAAC7aGIpFGWLJck7BoMAAADsHo7JnwAAAOzOqcWMvoEAAADsCtFMQ+FMnZsACQAAwG5i4kRhZUpy7AVAAAAAdlkFML0QAABwelLIt43FLuC9uAdge+MrEcD5m/+DSPJevt+Td9wlvn0nbycFgWxcYOtmAgAAvnNu0NPR//f7Wv3sx+WHfQZk287/mWxcUP0H3qbC4l4GhAAAgO/k9N9Izmnths+Mloc9Y7JtC0CySUG1F71aIgAIAADYCjZORgGA7RsAXkG1JlsoMhYEAABs8WoAtvkh4hgRAACwtaeX8plTfgnAiEsB53wqz/8xhkfzQAAAuLATkk0KMjaQdxmrAee8t7x8lsqnQ/nBpkwUy4QR4wICAMD5nIycoua8Gm/4CYXVmlya8qj4uT77907ZxikNjzyg7p23aPMrt2t4/KiCcoXhIQAA4HwFgBTV6pp+xRvzmwFxXg2PH1X3rn/W0T/9XW388z9JYSjDZYFdiyMP4Lxyg02lq8cZiAsgmmlo8lkv1PKv/K6mX/p6Kct4IoMAAIDztwrAqv+FD4HOz/8nTT77xfLDIRFAAAAAdougUFLzR35GcWeJbZkJAADAblJ63OWavPIFMpapgAAAAOwq1StfpKBY5jIAAQAA2E2SxYtkCkWxIRMBAADYTZNAqZzv+c/8TwAAAHYRL3ZjJAAAALuyADzPZhIAAACAAAAAADsT7wIAsG15lyldOabBoQOK6nOK5+a//uekQ3XvuEXpyRXFjbai2abC2iyDBxAAAMZJb/8+bd55i7p336rB4f1KTxzX4MgBzb7mLWq+6V9/3c93/Z6OfeQPdPILn1bSnFc401A4Oa2ks6TS4y5X6ZInKmq2GViAAABw4ZkzzuBT9Q/cq7UvXK31W27Q4MB9Gp44Lre5kW9R6zLJBvLp8Oy/krXyw4Gy1ePqDfrS/n35jweBgnJFYW1WyfweVZ/5fFWvfKHC2qxMEHAIQAAwBADO69QfRQrKE3LddZ266VqtfvpvdOpLN8itnZBzmSSTP5VmjEwYyjsrG0X5s+pnLwCZpCATxTJh9FBWSF7K1teUrq+pt3+fTv3jZ3T4Q7+j6hXfq9qLXqXSZU+WLZY5ICAAAOC8nPh76eT11+jk5z6h1c9+XEozmSjKJ3z7Dc7MH3WrWn+WTWxM/nsZc3qtwTun9MHDOn7Vh3Ti6r9U9Vkv1vRLX6uJy58pWyhxbEAAAMC5CwCr3v6vaP+vv0vZ5rpsnEhxfF5/fxMn8oOBVj7xEZ26/u9VfeYLNPPqH9HEk5/B8QEBAOxGrt+T63XlNrtymxvKNtblhwPZYkm2UFQ001RQmcx/cpZJXEc+68n4N9tT1m925b2XTYoX7s9prYLShLLuhlb+9iNa//I/avqlr1P9DT+hcGqG4wgCABjbeWg4VNbdkNtcV7axrmxzXW7jlLL1U8o21uS66/nH3Q1l3fX833tdqb8p1+/LDQdyg77ksnyiMkYmsDI2lEmKKizsUfHSJylutBRNN1TYc1m+jL3bT/Dzi/ffOAKs3TabzpkwlAknNDjygI588H3auP1mtX7851V+4hV8EYEAAC7sTO7lBn353qbS7ql8Et9YV7a+lk/g6yeVbawpW1+T21hXup5P7m5zQ663KT8cyKVD+XQon2XyaSplmXyWymejj707+/XlM/dH915yLv+5xqp7240yn75KJggUzTRVWL5EycJFqr3oXyhpL8nEya48XDaOZcJQ6o/ZnzspyLtM6zdeq/sO7FPjzT+l2Vf9sIxlrzQQAMAFka6t6oH3/YrWb/qCFATyw0E+kWfpaNL2o0fF3GiC9l9z9mkeebZpTv/P6UnemOBbfhmKOSMI/LAvP5D6B+5Vf/8+KYx0/G//TKVLn6TGG96m4mVPVFCa2G1rAOP7J7eBZAOlx4/p4G//qoZHDmrux37u0Z8+AAgA4BwuAKQDDQ4/oMHhA7LF4tdM3l/zgTGjJejzMVuMlrrP/O28U7a6orXrrtbadZ/W5HNeotnX/pgqT3mWtGvOJHfAO2WtlbzT0Q//T6UnT6j9k7+osDrFFyN2HNa3sP3PKG2Q/6VsbD7xPvRNZvudcRojE0YyQaATV/+17v2lt+vgH/yGBkcPcSjHazlAsoGOX/VBPfDbv6x07QRjAgIAwLcWAsFERa67riMf+l199dffpc27b2VcxqoBApko0conPqpD7/81ZRtrDAoIAADf4iQSJzJhqPV/+qy++pvv1satX2RQxun4BXkEHP/Yh3X4A78l199kUEAAAHgsZ5KxNu/6svb/9/do/UvXMyhjFgEKQx37iz/WsY/8EQMCAgDAY5lFjEwQqnffnTr8B7+hwZEDjMm4RZykox/+Pa1dfw0DAgIAwGMUBFq/9UYd+eDvyPW6jMdY/W1plXXXdeh//ZoGhwk4EAAAHtNCgJWc08on/1wnP/NxBmQMVwK699ymwx94b76pFEAAAPiWJ5Eoltvs6thff0iDg/czIONVcDI20Ml/+ITWrrua8QABAOAxfuHFibq336zVz34838kQ49MAYaR07YSOfvQPNTx2mAEBAQDgsXzlWSnLdOLvP6bBwfsYj7ELuII2bv6CTv7D3zIYIACAbcv7/OU/6VB+OMzfJzAcyA168oO+fJqOvg1HLwnK8ncLnOszyShS7/57tHn3bRyjcRME8s5p5eq/0uAYuzxiPPEuAOyCv6ytgmJJJgjzV9EGgUwQykSJfDpUtr6Wvw3Oe7nuuuRc/gbBXnf08+JztgrgBz2dvPaTqlz5AgXFMsdqjJi4oO6dX9apG67RzPe/iQEBAQBsq5N/55TMLWr29T+ueHZOkpONizJJQbZUlh8MlK0eP/363mzjlPygr+6+O9S7+1b17r1Lg6MHZZMkfyfBVk8iNlD3ji9peOyQgsWLOWDjFADWyPX7OnH1VZp87vcpnKwxKCAAgG2UAIpmGpp+6ese0xn2VDpUuraq3r13avWaj+nEp6+S721KwRZHgDFKT66od8/tKhAAY8fGiTZuu1HdW7+o6rNfzIBgvD5/GQLsgnM1+XT42P4fYaRouq7K056rzs/8shbf9euKFy6Sz7b42W9j5Ho9bdxxk9xwyKEau08tI9/f1Kkbr2UsQAAA228RwOffvt0vkqSgqee9Qsv/4bdVfsJT5fu9Lf7jeXXv+md5dgYczwYIIp364uc1PPwAgwECANiJihc9Ts03/7TC+pzcsL91E4ikbPWE0pVjDPKYrgIMDt2vza/czliAAAB2qsrTv1czr/yh/IkCl23Zr+t6XQ2OHmSAxzQAvHPq7iMAQAAAO/fv+iDQ9Iteq8LSJfJbdc3eGLl+T8PjRxngMeVdpo1bblC2ucFggAAAdqqkvajJ57xEplDcsm18Xb+n4QoBMLacV3//PqUPcgxBAAA7VxCo/PjLFZQrW7ZjoB8OlJ48wdiOq9EqzuAYl3EwPtgHAPh2vnBqdQXlitLV4wzGheYlnw3zrZzTC/MopXdZfh/HoQMcDxAAwM4OgFmF9Tn1eZ3vhZ//s1SyoYLJ6XxVRuZCFIDkvYasAIAAAHa2oFxRPNNU1wb5HgPGMCgX6lgUy2q88e2qvfCVsoXihQkASb7fU9xa4ICAAAB2POb8bVIA+T0ZevzljAXwGHATIPBtn/IxBAAIAAAAQAAAO5yRZLfwGoD3+Y1k+MZjBIAAAC7oXDQYKttY37JJyYShbFJkYB91wCUTxowDsIW4CRD4NmTrqxo+eFjeZzI2+s5LvFBS3GwzsI82+weBbFLY1qsBPsska2Us51UgAIAda7jyoNITx7fmRkDvZeNY4eQMA3sWxlj5LNPqNX+j/pH9Uppu2Q6MW/eHlNzmpspPerpqL341Bw0EALBT9e+/R9n6KclszdmeLZaVdJYZ2LMXgOQyrd90nU7deG0+2263RzBGe0HYSlU1jhgIAGBnytZOau3Ga+U2N2SC4DufOyQFE5MKZ+oM7jesJCuzTTdf8C5TUKyouOdxHCcQAMBO1du/T5t3fVneZTLhd/4lZCQVLrpMtlhmcL/ZSsC23X3JySaJ4jl2AsQYNTVDADyGv+YHfR2/6k81OHxANtqKu9K9TFJQ+QlXyEYJAzyuvGTLFUX1BmMBAgDYiVY+/mda/czH8rPRrdj/33kF1ZrKT/huthYeZ4FV6bInKapxGQcEALDjnPjUX+rQH/43ZZtdmSDakl/TZ5nKT3iqotk5BniMT/+tDVV63OUyEXsVYHxwDwDwTc/SnY7+39/X0Q+/X9nJE/nz6FvxHLr3MlGkqee+VLZYYpzHd/6XokDFix7PWIAAALYXI30bm7N457Rxyw06+uHf0/pN18kNBzJxsmWb0LhBXxNPerpKj+MtdmM9/3un8t7vUmEvAQACANhm87+Xsizfqe3M/faNkU9TuX4v373NOblBX+nJFa1/8XNavfbv1L//HqXrazKSTBht3Z8py2SiSNPf9wbFc/Mco3GWZao883kKq1OMBQgAYNvM/caqf2i/jvzx/5AtTsgN+pJ8vrd8FGt4/Ih6X92noFiS29xQ/4H75Ad9eefkvZOROSdbu7pBXxNPeaYqV3wPB2mcOadwclqVpz2XsQABAGyzAtDw+BEd/egHvv6/eS9jjGQf3szHWHv67n5jzs09sj4dypbKmn75Gzn7H/f5f9DX5NOeo8Le72IwQAAA23EVwMSJHrl97AV65s57+SzT1Pe+SLXnfz8HZ5xlqcLqpKZf9vot2hMCOL94DBC7KQXO+HaB5v8sVfGiy9T80Z+VLbHz37if/Zef/AxNPOFpDAYIAADfaMbIFExOa+5f/lsVuWN8rPl0qKBa08wr3yxbrjAgIAAAPMqE4ZxkQ9Vf9SZNPY+l/zE/mPJZqsnveRk3/4EAAPAN5ossldKBZl75JjV/9GcZkB1w9p90ltV47VvyTaGAMcVNgMC5nCyGQ8l4Tb/8jZp7y8/J8MKfMT/5d5KsZl/zoype9mQGBAQAgLOfKdpiSfU3vFWNH3q7gtIEgzLuskzV575Us6/6YcYCBACAs0/+0cycWm99p2oveY1MwJfaTjj7Txb3av5nf1m2yBMcIAAAnJ4hvLxzMkGgiac8S623/kL+ml+M/6HNUgWlCXXe8W7FjTYDAgIAwOi6cJZJgVHSWtT0y39A9df9OG/42wmMke/3ZeJEnXf8kqpXvpAxAQEA7Pqz/SyVzzLZKFY8v0fVZ71Q0y97Pc/47yCu35OR1H7bL2j6FT/IgIAAAHbtpJ8O5dKhbBgprM0qmVtQ5crnq/qM5/Na3512uIcD2ShW6yfeqfoPvJUBAQEAjOvZurxGuwCbM3YDNmd8Z05/6L3PN3xxTnKZfOZki0XF7SVFjZaKex+v6tOfp2Rhj5L2EmO80z5l0qGCypTa/+pdmnnVmxkQEADAOE7+tlTOJ2lr81f9Zpl8mkoPTe7OyftsdC0//zFbKiuqzSqYmlE001A0PaugVldx+RIlixcrLFdGLxjCI5nx/5RJU8WtRbXf/ouaet7LOaQgAIDxnP+dCnsu09K//00F5Up+Vu/zMJB38ml+HT+/nj/6lg5lbCAbxbLFsoKJqmySSIaNM7/peDuXr56MXQd4+TSTrFXliueq/Y5fUnHPpRxQEADAWJ+TGqugXFEwUWUwzjE3HEhZOlYrAflWzanCmYZmXvFGNX7oJ/lcAQEA7JBlAHmXMQ7bphKc3KAvE8cyNriwf45+T7ZU1sRTnq36D75V1We+gOMDAgAAHrMzbqZ8NNF0Q8HUtHr33yO3uS6bFCVzHlcMvJfrdWWTgspPfrpqL3mtpl/0agWVSY4fCAAAOFdn3XFnSQvv/C9a+8fP6vhVf6revjvknTvn91l4l8n1erJRrInvframnv/9mnz2SxTPdTguIAAA4FzLel2FtVnVX/sWVa/4Xq3f/Hmt/dPntH7TdcrWT0rWyhjz8ErCt7M64P3p7733knOKmh1Vn/ZslZ/8TFWe+mxF9RYHAwQAAJw3aSq3vqagXMn3UVjYo6kXv0a9r9yujdu/pO5tN6p7+80anngwjwCXyXsnOa98MwfpkZcZzvgxa2SslYJAxhhF9ZbKT7xCpcuepPJ3PVWF5UtkophjABAAAM47Y77urD4ollV+4hUqP/EKuf6bNVw5pv7996j/1X3qH7xfg2OHNDx2OF8heOimTp//WsZayVqFlSlFjZaiZkeFzrKSxYtVWNyrYLImy6QPEAAAtjebFJW0FpW0FqUrXyg5l1+/3+wq656SH/Tl+v18md9a2aQgE8cKylXZQpFXLwMEAICdUQRWxloFlUnu1Ae28kuLIQAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAHbr3zqW3fqAbYCvQgDnlbFW2eaG7GZXcqke+WKf3cc5l+90WCqf09chAwQAgAs4+xsNDn1VB3//1xQUy3LDwS4PACOf9hVMTKr1tncqmm7yOQICAMDODIDs1JpWr/mYvPNf+1LAXckPh4rbC2r+8E8xGCAAAOxg1soWSozDQwEQDhRMzciEEYOB8/ulyBAAwAXmfP56Y4AAAAAABAAAACAAAAAAAYAdj1vFAWCr8RQAtrcsVba+qqy7LnmXb5RijIwxks0/lrH5vz/aM2U8awYABADG7BN0uq6ld79Xg6MHla2tKj25omx1RenaitKTK0pXTyhdW1G2dlKuuy7v3Oj/Obqj2jm5zQ0pyxhMACAAMC5MGKm49/Eq7n38wz/oMnnn8sneOXmXSc7lj1F5p6y7oezUqoYnT2h47JCCpCBbLDOYAEAAYKzZQMYGj3pnQFCtSXPzKjJSGJvSNVyqAgEAYBdwTt577u+U5LNMGvQkx2UqEAAAdvSM56UgkLVW3nmdvl9jt4qkbLMrn6Z8boAAALBzJ/94rqOZV/+Yoqlp+eFg1w+JS4cKSmWFMw0+P0AAANi5ARDW6pp5xRsVlLgxcxcccIZgG2MjIADnVdbflNtYYyAAAgDAbmK8ePPdbjj3Hw7EXZ4EAABgF0lXVzQ8fkQKucpMAAAAdo3uHV9S9+5bZaOYwSAAAAC7Qba+qgev+pDS1RWZMGJACAAAwHjxj/leDT/s68j/eb/Wrvu0giJ7cW53XKABcAHmFm4C3O5MGCkoTnzLP3/wwP06+ud/pAf/+oMygc3f3AkCAAAenlgCBeUKA7HdGy1z6h36quzEhPxg8PU39Bsrea/BkQfUvevLOvGpv9DmPXfIcOMfAQAAZ59ZpPTUqmw6POP1zbt1LJwko7AyJRNvoxvmjFF66oQOvv+/Kpyoyp1lx0ZjRwFw9JCGxw7JRBHX/AkAAHgU1qr/wH267z//vGxSyF/jvKvn/0x+0Ffrbe9S5WnP2V5/OOfUv/9u9ZzLz/4f5aqNCUPZYonPbQIAAL6xrLuhjS9dz30Akrx38ulQ9ZPHt+Wfz0QxW/kQAACwRZOKtTKcMeYBkKWy0aSC0gSDgfOO2zQB4MLlkLz3rIaAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAEAAAsC14GWMkw357IAAAYBcx8mkqn6UMBc47tgIGcP4nvSwdvQhod++A57NM3nv5LOPTAgQAgJ084znF7SW13vJvFNbq8ukuP/N1eQAUL30inxsgAADs5ACQgnJFE1d8j8JqjfEALiDuAQBwfhsgTeV7mwwEQAAAAAACAAAAEAAAAIAAAAAABAAAACAAAACPDVsREwAAgF1ZAPl3npEgAAAAu4bvb8r1eqwCEAAAgN2kf/Crcr0uA0EAAAB2k1PXXyO3ucEKAAEAANgtBkcO6uR1V/M6YgIAALBb+EFfh//3b6l3/z0yNmBAdiHeBgjg/HPccn5Bh797Skf+5H1a+eRH8x9g+Z8AAIBzzSYFhTMNBuJCnPWnqTZuv0nH//JPdOKaj0neyYQRA0MAAMA5ZozS1eM68am/UDAxKe8yxuQ8jLmc03DlmLq33az1W67X8NhhyRqZMBYbABAAAHBeJqPB0YM68L7/KGO5Ben8nfpLPhvKDQaSczJhOFr2Z/InAADgvE1GXr7fY+o57/GlPLoCbvgDAQDgQuHsH7jwX4YMAQAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAACAAAAAAAQAAAAEAAAAIAAAAAABAAAACACcc957mSBkIAAAW47ZZTvXWZJo8yu368TVf6W4taiwWlMwUVEwUZUJIwYIAEAA7MgAKJTUvevLuu9Xf1rBZE1xva242VE8N6+ktaC4vah4tqVwtqGgMiWbFBg0AAABMPa8lwkCmaAkt9lV7947tXnPrfLDoWStguqU4tmWoulZhdN1xXMdxa1Fxe0Fxc15xbNzsoWiZLjSAwAgAMaSsYFkA5kolgp5HPjBQL39X1Hvvjvls0wmCBRMVPNv1SmF1Vq+WtBZVmFhr6L2guLphmy5kocBAIAAwLgVgRmtDgSSktMrBm4wkHvwiAZHD0kukwljmSSRTYqySXEUBW0li3sVzy0o6Swrmm0qmm0qKFcYVwAgADCuUaAgkHnox7yX0lTZcE3ZqVUNHzykzfvukLnhM/nPtYHCyZoK8xcprLdUWNyjuNFR3F5S0uzIFEsyYTQKDQAAAYCxiQIZMwqCMyZx7+XTVFKq4bFDGh49JC8vOS8ThopmGopqs4qanXylYH5JSWeP4mZHtliSLZa56RAACACMZRg8/C8P/XN6Z4jBscMaHDko3XaTvHOySUFBuaKgMqmo0VLSXlLcXlRx6RKFtVmFtRkFE5MKJqqMLQAQABjbPrBWslZSmIeB90pPrSo9uaLe/XfrlCQTxnkYVKqKZucUN9qK5/NVgrjRylcRGm2F1RoDCgAEAMZ1xcAEoRQofxJhdLnADXpyx7oaHD6gDWNkbCgThgoqkwqrNUWNlqLZOSVz84rm5pW0FxU35xXV5xhTACAAMIZFIFkjI/twFHgvyct7p3T1uNKVB9W77+788kIYyRSKCipVBaWKoum6koWLlMzvUdJeUrKwR2FtVjZKZEI+LQGAAMBYrRJIJv/OPvKGQ8nLdTfkuusaeqm3706t33SdFIayYSQTJ4qbHSULF6mwsFdJZ1lxZ1Hh5LSCckW2WOZJBAAgADCWYXC2+TtN5dJU6m1q8+QJde/8siQveSmYqChutBS39yjuLKrQXlTU7CiaaSqcmlEwUeVJBAAgADC+cfDwx2c+m+A2u9q89y5t3nO7vMvks0zBRPX0PQRxs62o0VEyN5/feNhsK5icVlCekB7xKwEACACMVRzkNx6GD0/nWab+Q9sfp6m894qqNYVT0wqnZxXWZhXXW4o7y0oWL1Y808xfllSuyMYJYwoABADGkrUyNh49iTBaKUiHGhw9qP6h/fJZKmMD2VI537egXFFQqSpuLqiw5xJFzfxNitFMQ9FMk3ciACAAgLFdKHho++No9APey2ep0tXjGq4ck5xT945bZD6fh4NNCgomJpXMLymut5V0lhS3lxTPzSuqtxRWpxhUAAQAMH5FYGTM6O2JejgKlGVyWVduc0Pp6nH1H7hXxuQbHpkwUlSbVTg1q6jZUtJezgOhs0eF+SXZuCgTxTJRxPgCIACAcYqCr73J8DTn5Ad9DY4c0ODwAfnbb5K8l40iBZVJBZVJxfW24s6SkvllFZYvVTQ7l19eKJVlSxOMLwACABjLOMg/ePidCJK8c0pXVzQ8kW9o5K/PZOJEQZTIVibzrY9biyosLCuaW1Ayt6BgalrR1LSCyhTjCoAAAMZ71WC026Ekea9s0FN2rKvB4f3SLdfLmPwSQjBZUzTTUFxvKWq2lbSWFDXb+b83WopmmowpAAIAGNcoOH1/QRiNdjlUvv3xyjGlDx5R945b8p8ahgomqgqqUwonawqnZlWYv0hJZ1FxZ1lxa1HRTEMmLjzy5YwAQAAAY7BKIJ2OAunhlQJ5r2z9lLJTJ9U/cK/kpVNRLJMUZJOCbBQrmptX3FlWceni/B6D9qLCqVkFpQkeUQRAAABjGQbG5K9Vlh7xNILr9+R6m5J3Gjx4WN3bbtSqDfJLCXGsqNFSYfkSJa1FJZ389crhdF1RbTZ/JwJPIgAgAIDxCwNzer3/kS9M8i6V30zVv+8e9e69O19FcPmNh8noPQhJe0lxa0Fxo62o3lLcaMmWKwqKJclYxhcAAQCM5arBQ5Fw5peo9+rt36fefXdrLcvkXf5OhHBqRnFjTtFsS/FcR3FzXoXFixXWZhXWZhSUq6wWACAAgLFeNTjLOxGGxw5p8ND2x8bKlkoKKjWFlUmF07OK620V9lyquNlR1Ggrmq4rrM1ybwEAAgAYW9bKjHYwlDS6hOCUrhzV8MHD0j23SkEgm5RkkyR/GmFiUsn8suK5BcWt+VEYdBTXWwomKowpQAAAGMtVgrNsf+xdqqw7VLa+Jvn92tx3u4wN8ncixLGCqRlF0w3F9TnFrUUl83uUdPKbD+1ENb9fgWcUAQIAwBhGgfTw/YYP7Vsw6Cvt95SeWlP/wL0yGj21EI9WCypTimabKszvybc/XrpYSXtJtlwZxUNy9jAgFgACAMD2jILTcXCW/+wH/XxDo+NH1bv3Tp264TMyUZTvW1CeUNJaVDy3oMLiXiXzexS3FvJgKFfy7wvF0e/hGWuAAAAw1oHgnFx3Q9nGuoaHHpD8taM3KIYKJiYVNdqKmx0VlvZqcPCrknM8kgiMw5f7zS+5hFQH8NiMdjr0zslnaf5tOJSJIgW8HRFgBQDADl4pGO10aEL+GgHGEet0AAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAEAAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAEAAAAIAAAAAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAAQAAAAEAAAAIAAAAAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAEAAAABAAAACAAAAAAAQAAAAgAAAAAAEAAAAIAAAAQAAAAAACAAAAEAAAAIAAAAAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAAQAAAAgAAAAIAAAAAABAAAACAAAAAAAQAAAAgAAABAAAAAAAIAAAAQAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAAEAAAAAAAgAAABAAAACAAAAAgAAAAAAEAAAAIAAAAAABAAAACAAAADBG/v8Att1dUuWoLq8AAAAASUVORK5CYII='

definition = ActionDefinition(
                name= slug,
                url= f'{action_hub}/actions/{slug}/action',
                label= 'PowerPoint',
                icon_data_uri= icon_data_uri,
                form_url= f'{action_hub}/actions/{slug}/form',
                supported_action_types= ['query'],
                description= 'This action will generate a PowerPoint slide based on a Look',
                params= [],
                supported_formats= ['json'],
                supported_formattings= ['unformatted'],
                supported_visualization_formattings= ['apply'],
            )

@app.post(f'/actions/{slug}/form')
def form():
    """Dummy endpoint for demonstration purposes."""
    return [
        ActionFormField(
            name='filename',
            label='Filename',
            description='Filename for the generated pptx document',
            required=True,
        ),
            ActionFormField(
            name='title',
            label='Title',
            description='Filename for the generated slide',
            # required=True,
        )
    ]

@app.post(f'/actions/{slug}/action')
def action(payload: ActionRequest):

    # create presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    # title.text = 'Test'
    title.text = payload.form_params['title']

    # Pull back data and format into series
    data = json.loads(payload.attachment.data)
    print(data[0])

    ##### Pull back Vis config Info########
    visconfig = payload.scheduled_plan.query['vis_config']
    # print(visconfig)

    # Get type of visualization (Column, pie, line, etc)
    vistype = visconfig['type']

    # numberformat = visconfig['label_value_format']

    # Get colors from Vis
    colors = visconfig['series_colors']
  
    # valuelabels
    valuelabels = visconfig['show_value_labels']


    # Initialize chart
    chart_data = ChartData()
    x, y, cx, cy = Inches(1.5), Inches(2.5), Inches(7), Inches(4.5)
  
    pivotsraw = payload.scheduled_plan.query['pivots']
    first_row = data[0] 
    dimension = list(first_row.keys())[0]
    measure = list(first_row.keys())[1]
    categories=[]

    # Map Data to Series Values
    if pivotsraw is not None:
        for i in data:
            categories.append(i[dimension])
        chart_data.categories = categories
        # maxvaluedefault = 0
        for pivot in first_row[measure].keys():
            for key in first_row[measure][pivot].keys():
                for dim in categories:
                    values = []
                    for entry in data:
                        # int(a.replace(',', ''))
                        # values.append(float(entry[measure][pivot][key].replace(',', '')))
                        values.append(float(entry[measure][pivot][key]))
                        # if float(entry[measure][pivot][key])>maxvaluedefault:
                        #     maxvaluedefault == float(entry[measure][pivot][key])
                setvalues = tuple(values)
                chart_data.add_series(key, setvalues)
    else:
        values = []
        for key in data:
            categories.append(key[dimension])
        for value in data:
            values.append(value[measure]) 
        setvalues = tuple(values)
        chart_data.categories = categories
        chart_data.add_series(dimension, setvalues)

  

    #yaxis
    # maxvalue = int(visconfig['y_axes'][0]['maxValue'])
    # if maxvalue=="":
    #     maxvalue == maxvaluedefault

 
    # Create Chart
    if vistype == 'looker_column':
        chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        series = chart.plots[0].series[0]
        value_axis = chart.value_axis
        # value_axis.maximum_scale = maxvalue
        series = chart.plots[0].series

        #totals
        if valuelabels:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            # data_labels.number_format = numberformat
            data_labels.font.size = Pt(13)
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        color_list = []
        for hex in colors:
            color_list.append(colors[hex])
       
        i = 0
        print(pivotsraw)
        for (bar) in series:
            fill = bar.format.fill
            fill.solid()
            if pivotsraw=="" or pivotsraw is None: 
                color = color_list[0]
                rgbcolors = ImageColor.getrgb(color)
                rgbcolors = tuple(rgbcolors)
                r,g,b = rgbcolors
                fill.fore_color.rgb = RGBColor(r, g, b)
            else:
                color = color_list[i]
                rgbcolors = ImageColor.getrgb(color)
                rgbcolors = tuple(rgbcolors)
                r,g,b = rgbcolors
                fill.fore_color.rgb = RGBColor(r, g, b)
                i+=1
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

    # elif vistype == 'looker_pie':
    #     chart = slide.shapes.add_chart(
    #         XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    #     ).chart

    elif vistype == 'looker_line':
        chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        ).chart
        series = chart.plots[0].series[0]
        line = series.format.line
        value_axis = chart.value_axis
        # value_axis.maximum_scale = maxvalue
        color_list = []
        for hex in colors:
            color_list.append(colors[hex])
        ## looks like line isn't iterable in this version? Need to dig more. I'm assuming one series for now
        color = color_list[0]
        rgbcolors = ImageColor.getrgb(color)
        rgbcolors = tuple(rgbcolors)
        r,g,b = rgbcolors
        line.color.rgb = RGBColor(r, g, b)
        
    else:
        slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        )

   
    
    #totals
    # if valuelabels:
    #     plot = chart.plots[0]
    #     plot.has_data_labels = True
    #     data_labels = plot.data_labels
    #     data_labels.number_format = numberformat
    #     data_labels.font.size = Pt(13)
    #     data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    # # Save to PPTX file
    prs.save(payload.form_params['filename'])
    # prs.save('test.pptx')

